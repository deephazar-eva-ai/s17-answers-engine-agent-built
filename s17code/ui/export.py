"""Turn a rendered answer into a downloadable DOCX or PPTX.

A pure transform over already-grounded text the client already has by the
time a run finishes: the question, the answer text (markdown headers/bullets
over prose), and the source list. This module makes no LLM call and touches
no run state.

Parsing mirrors the client's own md()/mdInline() in answers.html (headers,
bullets, **bold**/*italic*) so the exported file matches what the page
rendered.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field

from docx import Document
from docx.oxml.ns import qn
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt as PptxPt

_HEADER_RE = re.compile(r"^(#{1,4})\s+(.*)")
_BULLET_RE = re.compile(r"^[*\-]\s+(.*)")
_INLINE_RE = re.compile(r"(\*\*([^*]+)\*\*|\*([^*]+)\*|_([^_]+)_)")


@dataclass
class Block:
    kind: str  # "h" | "p" | "ul"
    text: str = ""
    level: int = 1
    items: list[str] = field(default_factory=list)


def parse_markdown(text: str) -> list[Block]:
    """The same three shapes the client renders: headers, paragraphs, bullet lists."""
    blocks: list[Block] = []
    ul: Block | None = None
    for raw in str(text or "").splitlines():
        line = raw.strip()
        if not line:
            ul = None
            continue
        h = _HEADER_RE.match(line)
        b = _BULLET_RE.match(line)
        if h:
            ul = None
            blocks.append(Block("h", text=h.group(2), level=len(h.group(1))))
        elif b:
            if ul is None:
                ul = Block("ul")
                blocks.append(ul)
            ul.items.append(b.group(1))
        else:
            ul = None
            blocks.append(Block("p", text=line))
    return blocks


def _inline_runs(text: str) -> list[tuple[str, bool, bool]]:
    """Split text into (text, bold, italic) runs on **bold**/*italic*/_italic_."""
    runs: list[tuple[str, bool, bool]] = []
    last = 0
    for m in _INLINE_RE.finditer(text):
        if m.start() > last:
            runs.append((text[last : m.start()], False, False))
        if m.group(2) is not None:
            runs.append((m.group(2), True, False))
        else:
            runs.append((m.group(3) if m.group(3) is not None else m.group(4), False, True))
        last = m.end()
    if last < len(text):
        runs.append((text[last:], False, False))
    return runs or [(text, False, False)]


def slug(text: str, max_len: int = 60) -> str:
    s = re.sub(r"[^a-zA-Z0-9]+", "-", str(text or "")).strip("-").lower()
    return (s[:max_len].rstrip("-")) or "answer"


def _host(url: str) -> str:
    m = re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://([^/]+)", url)
    return m.group(1) if m else url


def _add_hyperlink_docx(paragraph, url: str, text: str) -> None:
    part = paragraph.part
    r_id = part.relate_to(
        url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = paragraph._p.makeelement(qn("w:hyperlink"), {qn("r:id"): r_id})
    run = paragraph._p.makeelement(qn("w:r"), {})
    rpr = paragraph._p.makeelement(qn("w:rPr"), {})
    color = paragraph._p.makeelement(qn("w:color"), {qn("w:val"): "1D4ED8"})
    underline = paragraph._p.makeelement(qn("w:u"), {qn("w:val"): "single"})
    rpr.append(color)
    rpr.append(underline)
    run.append(rpr)
    t = paragraph._p.makeelement(qn("w:t"), {})
    t.text = text
    run.append(t)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def build_docx(question: str, answer: str, sources: list[tuple[str, str]]) -> bytes:
    doc = Document()
    doc.add_heading(question or "Answer", level=0)

    for block in parse_markdown(answer):
        if block.kind == "h":
            doc.add_heading(block.text, level=min(block.level + 1, 4))
        elif block.kind == "ul":
            for item in block.items:
                p = doc.add_paragraph(style="List Bullet")
                for txt, bold, italic in _inline_runs(item):
                    run = p.add_run(txt)
                    run.bold = bold
                    run.italic = italic
        else:
            p = doc.add_paragraph()
            for txt, bold, italic in _inline_runs(block.text):
                run = p.add_run(txt)
                run.bold = bold
                run.italic = italic

    if sources:
        doc.add_heading("Sources", level=1)
        for url, title in sources:
            p = doc.add_paragraph(style="List Number")
            p.add_run(f"{title or url} ").bold = False
            _add_hyperlink_docx(p, url, url)

    for p in doc.paragraphs:
        for run in p.runs:
            run.font.size = run.font.size or Pt(11)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


_MAX_BULLET_CHARS = 220


def _slide_bullets(prs: Presentation, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:120]
    body = slide.placeholders[1].text_frame
    body.clear()
    first = True
    for bullet in bullets or ["(no detail)"]:
        text = bullet if len(bullet) <= _MAX_BULLET_CHARS else bullet[: _MAX_BULLET_CHARS - 1] + "\u2026"
        para = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        para.level = 0
        for txt, bold, italic in _inline_runs(text):
            run = para.add_run()
            run.text = txt
            run.font.bold = bold
            run.font.italic = italic
            run.font.size = PptxPt(18)
    return slide


def build_pptx(question: str, answer: str, sources: list[tuple[str, str]]) -> bytes:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = (question or "Answer")[:150]
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Answers engine"

    blocks = parse_markdown(answer)
    current_title = "Answer"
    current_bullets: list[str] = []
    made_slide = False

    def flush():
        nonlocal made_slide
        if current_bullets:
            _slide_bullets(prs, current_title, current_bullets)
            made_slide = True

    for block in blocks:
        if block.kind == "h":
            flush()
            current_title = block.text
            current_bullets = []
        elif block.kind == "ul":
            current_bullets.extend(block.items)
        else:
            current_bullets.append(block.text)
    flush()

    if not made_slide:
        _slide_bullets(prs, "Answer", [answer.strip() or "No answer text."])

    if sources:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Sources"
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, (url, title) in enumerate(sources[:12]):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = f"{title or _host(url)} \u2014 {_host(url)}"
            para.level = 0
            run = para.runs[0] if para.runs else para.add_run()
            run.hyperlink.address = url
            for r in para.runs:
                r.font.size = PptxPt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()